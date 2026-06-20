"""Generate executive PPT for Oracle AI DBA Assistant. Requires: pip install python-pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_bullet_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide.shapes.title.text = title
    # Left box
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.3), Inches(5))
    tf = left.text_frame
    tf.text = left_title
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(16)
    for item in left_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
    # Right box
    right = slide.shapes.add_textbox(Inches(5.0), Inches(1.5), Inches(4.3), Inches(5))
    tf = right.text_frame
    tf.text = right_title
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(16)
    for item in right_items:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
    return slide

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

add_title_slide(
    prs,
    "Oracle AI DBA Assistant",
    "Governed, SOP-Driven Database Operations Platform\n[Organization] | [Date]"
)

add_bullet_slide(prs, "Executive Summary", [
    "Problem: Manual SOPs, ad-hoc SQL, inconsistent incident response",
    "Solution: AI assistant + SOP playbooks + live diagnostics + human approval",
    "Current: Working prototype with SOP matching, diagnostics, and change requests",
    "Value: Faster ops, lower risk, repeatable compliance-friendly processes",
    "Next: Multi-environment, audit logging, production-grade governance",
], "Emphasize human-in-the-loop, not autonomous AI.")

add_bullet_slide(prs, "Why We Need This", [
    "SOPs live in documents — slow during incidents",
    "Manual SQL — errors and inconsistency",
    "Tribal knowledge — key-person dependency",
    "Limited audit trail — compliance gaps",
    "Single-environment — hard to scale safely",
])

add_two_column_slide(prs, "What the Platform Does",
    "User asks in plain English", [
        "Who is blocking my session?",
        "Kill stuck session",
        "Add 10GB to USERS tablespace",
    ],
    "System delivers", [
        "Matches SOP and runs diagnostic SQL",
        "Shows checklist + action preview",
        "Plans change + policy check",
        "DBA approves before execution",
    ])

add_bullet_slide(prs, "Platform Architecture", [
    "Layer 1: Business Users — DBAs, Support, Change Management",
    "Layer 2: AI Assistant UI — chat, SOP checklists, approvals",
    "Layer 3: Intelligent Routing — change requests, SOPs, known tasks",
    "Layer 4: Governance — policy, approval gate, action guard ★",
    "Layer 5: MCP Server → Oracle Database (DEV / UAT / PROD)",
], "Highlight governance layer as the trust anchor.")

add_bullet_slide(prs, "SOP Workflow — Controlled Execution", [
    "1. Match — question matched to SOP playbook",
    "2. Diagnose — live Oracle query (sessions, locks, tablespace)",
    "3. Review — DBA reads checklist and results",
    "4. Approve — human selects target and approves",
    "5. Execute & Verify — action runs; outcome confirmed",
])

add_two_column_slide(prs, "Risk Management & Controls",
    "In place today", [
        "Human approval before execution",
        "SOP diagnostic SQL (evidence first)",
        "Change-request policy engine",
        "Read-only known tasks",
        "UI / execution server separation",
    ],
    "Planned (recommended)", [
        "Environment-based policy (DEV/UAT/PROD)",
        "Full audit log",
        "SSO / RBAC",
        "Post-action verification",
        "Dual approval for high-risk PROD actions",
    ])

add_bullet_slide(prs, "Roadmap & Recommendation", [
    "Phase 1 (4–6 wks): Stabilize — routing, config, DB identity",
    "Phase 2 (6–8 wks): Govern — policy, audit, verification, tests",
    "Phase 3 (ongoing): Scale — multi-DB, more SOPs, ticket integration",
    "Recommendation: Approve Phase 1 + 2 for UAT/PROD pilot",
    "Ask: Budget, DBA SOP sponsor, multi-env connectivity, security review",
])

out = "Oracle_AI_DBA_Assistant_Executive_Deck.pptx"
prs.save(out)
print(f"Created: {out}")